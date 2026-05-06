#!/usr/bin/env python3
"""
generate_ingestion_deck.py  —  v4
17 executive slides: Full lifecycle with timings · Architecture diagram
· ROI calculation · Hero numbers · NR branding

Usage:  python3 generate_ingestion_deck.py
Output: Ingestion-Team-AI-Led-Demo.pptx  +  .pdf
"""

import os, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
G      = RGBColor(0x00, 0xAC, 0x69)   # NR green
GB     = RGBColor(0x00, 0xD0, 0x84)   # NR bright
GD     = RGBColor(0x00, 0x6B, 0x40)   # dark green
NVY    = RGBColor(0x0D, 0x0D, 0x24)   # deep navy
CHR    = RGBColor(0x1F, 0x29, 0x37)   # charcoal
W      = RGBColor(0xFF, 0xFF, 0xFF)
LG     = RGBColor(0xF3, 0xF4, 0xF6)   # light grey bg
MG     = RGBColor(0xE5, 0xE7, 0xEB)   # mid grey border
GR     = RGBColor(0x6B, 0x72, 0x80)   # body grey
BK     = RGBColor(0x14, 0x14, 0x1E)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
AMB    = RGBColor(0xFF, 0xC1, 0x07)   # amber
BLU    = RGBColor(0x00, 0x52, 0xCC)   # Atlassian
PUR    = RGBColor(0x7C, 0x3A, 0xED)
TEAL   = RGBColor(0x06, 0x82, 0xA0)   # teal for docs

FTR = "© 2025 New Relic, Inc.  All rights reserved.  Confidential and proprietary.  For internal use only."
DIR = os.path.dirname(os.path.abspath(__file__))
ARCH_IMG = os.path.join(DIR, "architecture (1).png")

SW, SH = Inches(13.33), Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def prs():
    p = Presentation(); p.slide_width = SW; p.slide_height = SH; return p

def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])

def R(s, x, y, w, h, fill=None, ln=None, lw=Pt(0.75)):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = ln if ln else (MG if ln is None else None)
    if ln:   sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def T(s, t, x, y, w, h, sz=12, bold=False, c=BK,
      a=PP_ALIGN.LEFT, it=False, f="Calibri", wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = t
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = it
    r.font.color.rgb = c; r.font.name = f
    return tb

def HL(s, x, y, w, c=MG):
    R(s, x, y, w, Inches(0.018), fill=c)

def ftr(s):
    T(s, FTR, Inches(0.3), Inches(7.06), Inches(11.6), Inches(0.36), sz=7, c=GR)
    T(s, "⬢ new relic", Inches(11.8), Inches(7.01), Inches(1.4), Inches(0.42),
      sz=10, bold=True, c=G, a=PP_ALIGN.RIGHT)

def hdr(s, title, h=Inches(0.68)):
    R(s, 0, 0, SW, h, fill=G)
    T(s, title, Inches(0.4), Inches(0.07), Inches(12.5), h - Inches(0.07),
      sz=23, bold=True, c=W, a=PP_ALIGN.CENTER)

def stripe(s, x, y, w, h, color):
    """Thin left accent stripe on a card."""
    R(s, x, y, Inches(0.07), h, fill=color)

def card(s, x, y, w, h, fill=W, border=MG):
    R(s, x, y, w, h, fill=fill, ln=border, lw=Pt(0.85))

def top_bar(s, x, y, w, h=Inches(0.06), color=G):
    R(s, x, y, w, h, fill=color)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def s01_title(p):
    s = blank(p)
    R(s, 0, 0, SW, SH, fill=NVY)
    R(s, 0, 0, SW, Inches(0.06), fill=G)
    R(s, 0, SH - Inches(0.06), SW, Inches(0.06), fill=G)
    R(s, Inches(8.6), 0, Inches(4.73), SH, fill=RGBColor(0x12, 0x12, 0x30))
    T(s, "⬢ new relic", Inches(0.6), Inches(0.32), Inches(3), Inches(0.52),
      sz=22, bold=True, c=G)
    T(s, "Ingestion Team",           Inches(0.6), Inches(1.52), Inches(7.6), Inches(0.72), sz=38, bold=True, c=W)
    T(s, "AI-Led Development",       Inches(0.6), Inches(2.26), Inches(7.6), Inches(0.65), sz=34, bold=True, c=G)
    T(s, "Environment",              Inches(0.6), Inches(2.93), Inches(7.6), Inches(0.62), sz=34, bold=True, c=W)
    HL(s, Inches(0.6), Inches(3.7), Inches(6.8), G)
    T(s, "Claude Code  +  GitHub Copilot  +  New Relic MCP",
      Inches(0.6), Inches(3.86), Inches(7.6), Inches(0.42), sz=15, c=GR, it=True)

    chips = [
        ("🚀", "10×  faster feature delivery"),
        ("🧠", "Onboard in hours, not weeks"),
        ("🔍", "30-min incident diagnosis"),
        ("⚡", "5 agents working in parallel"),
        ("📦", "Knowledge that compounds"),
    ]
    for i, (ic, lb) in enumerate(chips):
        yy = Inches(1.2 + i * 1.05)
        R(s, Inches(8.9), yy, Inches(4.0), Inches(0.82),
          fill=RGBColor(0x1A, 0x1A, 0x3A), ln=RGBColor(0x2A, 0x2A, 0x5A))
        R(s, Inches(8.9), yy, Inches(0.06), Inches(0.82), fill=G)
        T(s, ic, Inches(9.06), yy + Inches(0.2), Inches(0.4), Inches(0.42), sz=18)
        T(s, lb, Inches(9.52), yy + Inches(0.2), Inches(3.25), Inches(0.42), sz=13, c=W)


def s02_problem(p):
    s = blank(p)
    hdr(s, "The Challenge — Four Real Pains the DataOS Ingestion Team Was Living With")

    cards = [
        ("🕐", "New Engineer Onboarding",   "2+ WEEKS",   "to reach full productivity in the codebase",  RED),
        ("🔍", "Incident Root Cause",        "1–10 HOURS", "backtracking pipelines one by one manually",  AMB),
        ("💥", "Hidden Blast Radius",        "5-DAY",      "incident from what looked like a '1-day fix'", CHR),
        ("🗂", "Context After Absence",      "HOURS–DAYS", "to rebuild codebase understanding after a break", CHR),
    ]
    positions = [(Inches(0.4), Inches(0.82)), (Inches(6.88), Inches(0.82)),
                 (Inches(0.4), Inches(3.92)), (Inches(6.88), Inches(3.92))]

    for (ic, title, big, sub, accent), (x, y) in zip(cards, positions):
        w, h = Inches(6.05), Inches(2.88)
        card(s, x, y, w, h)
        top_bar(s, x, y, w, color=accent)
        T(s, ic + "  " + title, x + Inches(0.22), y + Inches(0.2), w - Inches(0.44), Inches(0.44), sz=15, bold=True)
        T(s, big,  x + Inches(0.22), y + Inches(0.72), w - Inches(0.44), Inches(0.9),
          sz=48, bold=True, c=accent, a=PP_ALIGN.CENTER)
        HL(s, x + Inches(0.22), y + Inches(1.72), w - Inches(0.44))
        T(s, sub,  x + Inches(0.22), y + Inches(1.82), w - Inches(0.44), Inches(0.42), sz=13, c=GR)

    # Total cost banner
    R(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.0))  # spacer
    ftr(s)


def s03_how_we_work(p):
    """DataOS Ingestion — how the team actually runs AI-led development."""
    s = blank(p)
    hdr(s, "DataOS Ingestion  —  How We Run AI-Led Development Today")

    # Left: what we own
    card(s, Inches(0.4), Inches(0.82), Inches(4.1), Inches(5.75))
    R(s, Inches(0.4), Inches(0.82), Inches(4.1), Inches(0.48), fill=CHR)
    T(s, "What We Own", Inches(0.6), Inches(0.9), Inches(3.7), Inches(0.34), sz=14, bold=True, c=W)
    repos = [
        ("om-airflow-dags",       "OpenMetadata ingestion pipelines\n(Airflow DAGs, connectors, plugins)"),
        ("spark-kafka-apps",      "Spark + Kafka streaming pipelines\n(real-time ingestion to warehouse)"),
        ("tf-dataos-new-relic",   "Terraform NR monitoring\n(alerts, dashboards — 55 commits/month)"),
    ]
    for i, (repo, desc) in enumerate(repos):
        y = Inches(1.48 + i * 1.48)
        R(s, Inches(0.55), y, Inches(3.8), Inches(1.3), fill=LG, ln=MG, lw=Pt(0.75))
        top_bar(s, Inches(0.55), y, Inches(3.8), color=G)
        T(s, repo, Inches(0.7), y + Inches(0.12), Inches(3.5), Inches(0.36),
          sz=11, bold=True, c=GD, f="Courier New")
        T(s, desc, Inches(0.7), y + Inches(0.52), Inches(3.5), Inches(0.68),
          sz=10.5, c=BK)

    # Middle: daily workflow
    card(s, Inches(4.72), Inches(0.82), Inches(4.1), Inches(5.75))
    R(s, Inches(4.72), Inches(0.82), Inches(4.1), Inches(0.48), fill=G)
    T(s, "Daily Workflow", Inches(4.92), Inches(0.9), Inches(3.7), Inches(0.34), sz=14, bold=True, c=W)
    steps = [
        ("1", "Open Claude Code in repo",    "Skills auto-load context + graph"),
        ("2", "GitHub Copilot in IDE",        "Line completion as we code"),
        ("3", "Work — skills auto-trigger",   "/airflow /pyspark /nrql /terraform…"),
        ("4", "NR MCP for incidents",         "Query logs, metrics, traces in-chat"),
        ("5", "/avengers for complex tasks",  "Parallel agents — 1 day → 2 hours"),
        ("6", "/wrap-up at session end",      "Jira updated · lessons captured · graph refreshed"),
    ]
    for i, (num, title, sub) in enumerate(steps):
        y = Inches(1.46 + i * 0.72)
        R(s, Inches(4.87), y, Inches(3.8), Inches(0.66),
          fill=(RGBColor(0xF0, 0xFB, 0xF5) if i % 2 == 0 else W), ln=MG, lw=Pt(0.5))
        R(s, Inches(4.87), y, Inches(0.36), Inches(0.66), fill=G)
        T(s, num,   Inches(4.87), y + Inches(0.17), Inches(0.36), Inches(0.34),
          sz=11, bold=True, c=W, a=PP_ALIGN.CENTER)
        T(s, title, Inches(5.32), y + Inches(0.04), Inches(3.2), Inches(0.32), sz=11, bold=True)
        T(s, sub,   Inches(5.32), y + Inches(0.36), Inches(3.2), Inches(0.28), sz=9.5, c=GR)

    # Right: tools stack
    card(s, Inches(9.05), Inches(0.82), Inches(3.85), Inches(5.75))
    R(s, Inches(9.05), Inches(0.82), Inches(3.85), Inches(0.48), fill=GD)
    T(s, "Tools in Use", Inches(9.25), Inches(0.9), Inches(3.45), Inches(0.34), sz=14, bold=True, c=W)
    tools = [
        ("Claude Code",        "Team environment\n+ 20 skills",        G),
        ("GitHub Copilot",     "Line-level AI\nin IDE",                 CHR),
        ("NR MCP",             "Live observability\n12 tools",          BLU),
        ("Jira + Confluence",  "Auto-updated\nevery session",           TEAL),
        ("Code Graph",         "Blast radius +\nimpact analysis",       PUR),
    ]
    for i, (name, desc, color) in enumerate(tools):
        y = Inches(1.46 + i * 0.98)
        R(s, Inches(9.2), y, Inches(3.52), Inches(0.88), fill=LG, ln=MG, lw=Pt(0.75))
        R(s, Inches(9.2), y, Inches(0.07), Inches(0.88), fill=color)
        T(s, name, Inches(9.36), y + Inches(0.06), Inches(2.0), Inches(0.38), sz=12, bold=True, c=color)
        T(s, desc, Inches(9.36), y + Inches(0.46), Inches(3.1), Inches(0.36), sz=10, c=GR)
    ftr(s)


def s03b_two_tools(p):
    """Copilot + Claude Code — how they work together."""
    s = blank(p)
    R(s, 0, 0, SW, Inches(1.52), fill=G)
    T(s, "Two AI Tools.  Different Jobs.  Better Together.",
      Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.95),
      sz=26, bold=True, c=W, a=PP_ALIGN.CENTER)

    # Copilot
    card(s, Inches(0.4), Inches(1.7), Inches(5.75), Inches(4.92), fill=LG)
    R(s, Inches(0.4), Inches(1.7), Inches(5.75), Inches(0.52), fill=CHR)
    T(s, "GitHub Copilot  —  Line-level AI",
      Inches(0.6), Inches(1.78), Inches(5.35), Inches(0.38), sz=16, bold=True, c=W)
    for i, (ic, txt_) in enumerate([
        ("✅", "Autocomplete lines & functions in IDE"),
        ("✅", "Inline chat: explain or refactor code"),
        ("✅", "Generate tests & boilerplate fast"),
        ("", ""),
        ("—", "No cross-session codebase memory"),
        ("—", "Can't query NR, Jira, Confluence"),
        ("—", "No impact / blast-radius analysis"),
        ("—", "No multi-agent parallel work"),
    ]):
        if not ic:
            HL(s, Inches(0.62), Inches(2.42 + i * 0.46), Inches(5.3))
            continue
        c_ = BK if ic == "✅" else GR
        T(s, ic + "  " + txt_, Inches(0.62), Inches(2.42 + i * 0.46), Inches(5.3), Inches(0.42), sz=12, c=c_)

    # Claude Code
    card(s, Inches(7.15), Inches(1.7), Inches(5.75), Inches(4.92),
         fill=RGBColor(0xF0, 0xFB, 0xF5), border=G)
    R(s, Inches(7.15), Inches(1.7), Inches(5.75), Inches(0.52), fill=G)
    T(s, "Claude Code  —  Team environment AI",
      Inches(7.35), Inches(1.78), Inches(5.35), Inches(0.38), sz=16, bold=True, c=W)
    for i, txt_ in enumerate([
        "✅  Full codebase context at every session start",
        "✅  Queries NR logs, metrics, alerts in real time",
        "✅  Knows what breaks before you write a line",
        "✅  5 parallel agents on one task with /avengers",
        "✅  Replayable golden patterns from past fixes",
        "✅  20 skills auto-trigger on file type or topic",
        "✅  Jira + Confluence auto-update each session",
        "✅  5-layer code review: CC + Copilot + agents + human",
    ]):
        T(s, txt_, Inches(7.35), Inches(2.42 + i * 0.46), Inches(5.35), Inches(0.42), sz=12, c=BK)

    T(s, "+", Inches(6.28), Inches(3.62), Inches(0.75), Inches(0.75),
      sz=42, bold=True, c=G, a=PP_ALIGN.CENTER)
    T(s, "Together:\n8–12×", Inches(5.9), Inches(4.44), Inches(1.52), Inches(0.72),
      sz=12, bold=True, c=GD, a=PP_ALIGN.CENTER)
    ftr(s)


def s04_architecture(p):
    """Full architecture diagram slide."""
    s = blank(p)
    hdr(s, "System Architecture  —  How It All Connects Under One Environment")

    if os.path.exists(ARCH_IMG):
        # Fill most of the slide with the architecture image
        s.shapes.add_picture(ARCH_IMG, Inches(0.3), Inches(0.78), width=Inches(12.73))
    else:
        T(s, f"[Architecture image not found at: {ARCH_IMG}]",
          Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.5), sz=14, c=RED, a=PP_ALIGN.CENTER)

    ftr(s)


def s05_lifecycle(p):
    """THE CENTERPIECE — lifecycle with time savings at each phase."""
    s = blank(p)
    hdr(s, "Full AI-Led Development Lifecycle  —  Time Savings at Every Phase")

    # Before/After times: verified by engineer (marked) or omitted if not confirmed
    # "—" = scope-dependent or not confirmed; no invented numbers
    phases = [
        ("01", "Requirement",  "Manager\n/ PM",      CHR,  W,    "—",      "—"),
        ("02", "Spec & Plan",  "AI-assisted",         G,    W,    "4–6h",   "20m–1h"),
        ("03", "Architect",    "AI + Graph",          GD,   W,    "2–4d",   "1–2d"),
        ("04", "Build",        "Task &\nvolume dep.", G,    W,    "2–3d",   "2–6h"),
        ("05", "Review",       "5-layer\nAI+Human",  GD,   W,    "—",      "—"),
        ("06", "Test & Gate",  "AI auto-\nmated",     G,    W,    "—",      "—"),
        ("07", "Deploy",       "Human\ninitiates",    CHR,  W,    "—",      "—"),
        ("08", "Monitor",      "NR MCP\n(AI)",        BLU,  W,    "1–10h",  "0–20m"),
        ("09", "Document",     "AI auto-\nupdate",    TEAL, W,    "manual", "auto"),
        ("10", "Learn",        "Golden\npatterns",    PUR,  W,    "—",      "auto"),
    ]

    BW = Inches(1.24)
    BH = Inches(1.32)
    GAP = Inches(0.01)
    total_w = len(phases) * BW + (len(phases) - 1) * GAP
    sx = (SW - total_w) / 2

    for i, (num, title, who, fill, tc, before, after) in enumerate(phases):
        x = sx + i * (BW + GAP)
        y = Inches(0.82)

        # Phase box
        R(s, x, y, BW, BH, fill=fill)
        T(s, num,   x + Inches(0.08), y + Inches(0.06), BW - Inches(0.16), Inches(0.24),
          sz=9, bold=True, c=RGBColor(0xAA, 0xFF, 0xAA) if fill != CHR else W)
        T(s, title, x + Inches(0.06), y + Inches(0.28), BW - Inches(0.12), Inches(0.52),
          sz=11, bold=True, c=tc, a=PP_ALIGN.CENTER)
        T(s, who,   x + Inches(0.06), y + Inches(0.84), BW - Inches(0.12), Inches(0.38),
          sz=8, c=RGBColor(0xCC, 0xFF, 0xCC) if fill not in (W, LG) else GR,
          a=PP_ALIGN.CENTER)

        # Arrow
        if i < len(phases) - 1:
            T(s, "›", x + BW - Inches(0.08), y + Inches(0.42), Inches(0.22), Inches(0.48),
              sz=16, bold=True, c=RGBColor(0xBB, 0xBB, 0xCC), a=PP_ALIGN.CENTER)

        # Time savings row
        by = y + BH + Inches(0.08)
        th = Inches(0.44)
        R(s, x, by, BW, th, fill=RGBColor(0xFF, 0xEE, 0xEE) if before not in ("∞", "0m") else LG)
        T(s, before, x + Inches(0.04), by + Inches(0.06), BW - Inches(0.08), Inches(0.3),
          sz=9, bold=True, c=RED if before not in ("∞", "0m") else GR, a=PP_ALIGN.CENTER)

        by2 = by + th + Inches(0.04)
        R(s, x, by2, BW, th, fill=RGBColor(0xF0, 0xFB, 0xF5))
        T(s, after, x + Inches(0.04), by2 + Inches(0.06), BW - Inches(0.08), Inches(0.3),
          sz=9, bold=True, c=GD, a=PP_ALIGN.CENTER)

    # Column labels on left
    label_y = Inches(0.82) + BH + Inches(0.08)
    T(s, "Before:", Inches(0.05), label_y + Inches(0.08), Inches(0.55), Inches(0.3), sz=8, c=RED, bold=True)
    T(s, "After:",  Inches(0.05), label_y + Inches(0.56), Inches(0.55), Inches(0.3), sz=8, c=GD,  bold=True)

    # Bottom summary — only confirmed phases shown; "—" = not yet measured
    R(s, Inches(0.4), Inches(5.96), Inches(12.5), Inches(0.62), fill=G)
    T(s, "Confirmed by engineer  ·  Spec+Arch+Build: days → hours  ·  Kafka diagnosis: hours → 0–20 min  ·  Overall: 8–12×",
      Inches(0.6), Inches(6.08), Inches(12.1), Inches(0.42),
      sz=15, bold=True, c=W, a=PP_ALIGN.CENTER)
    ftr(s)


def s06_onboard(p):
    s = blank(p)
    hdr(s, "Onboard in Hours, Not Weeks  —  Full Context on Day One")

    # Before
    card(s, Inches(0.4), Inches(0.82), Inches(5.8), Inches(3.72))
    top_bar(s, Inches(0.4), Inches(0.82), Inches(5.8), color=RED)
    T(s, "BEFORE", Inches(0.6), Inches(0.98), Inches(5.4), Inches(0.34), sz=12, bold=True, c=GR)
    T(s, "2 WEEKS", Inches(0.6), Inches(1.36), Inches(5.4), Inches(1.02),
      sz=64, bold=True, c=RED, a=PP_ALIGN.CENTER)
    T(s, "50 Slack DMs  ·  git archaeology  ·  tribal knowledge\nStill not fully productive by end of week 2",
      Inches(0.6), Inches(2.52), Inches(5.4), Inches(0.78), sz=12, c=GR, a=PP_ALIGN.CENTER)

    T(s, "→", Inches(6.38), Inches(2.0), Inches(0.55), Inches(0.7),
      sz=34, bold=True, c=G, a=PP_ALIGN.CENTER)

    # After
    card(s, Inches(7.1), Inches(0.82), Inches(5.8), Inches(3.72), border=G,
         fill=RGBColor(0xF0, 0xFB, 0xF5))
    top_bar(s, Inches(7.1), Inches(0.82), Inches(5.8), color=G)
    T(s, "AFTER", Inches(7.3), Inches(0.98), Inches(5.4), Inches(0.34), sz=12, bold=True, c=GD)
    T(s, "2 HOURS", Inches(7.3), Inches(1.36), Inches(5.4), Inches(1.02),
      sz=64, bold=True, c=G, a=PP_ALIGN.CENTER)
    T(s, "/bootstrap: detects stack, seeds memory, builds graph\nFull codebase context ready in 15 seconds",
      Inches(7.3), Inches(2.52), Inches(5.4), Inches(0.78), sz=12, c=GD, a=PP_ALIGN.CENTER)

    # Three outcome cards
    cs = [
        ("Full context\nDay 1",     "Architecture + lessons + todos\nauto-loaded every session"),
        ("Never re-explain\nthe codebase", "Team knowledge lives in\nmemory, not in one person"),
        ("4 weeks/year\nsaved",     "2 new engineers × 2 weeks\n= 4 engineering-weeks back"),
    ]
    for i, (title, desc) in enumerate(cs):
        x = Inches(0.4 + i * 4.3)
        card(s, x, Inches(4.72), Inches(4.06), Inches(1.92))
        top_bar(s, x, Inches(4.72), Inches(4.06), color=G)
        T(s, title, x + Inches(0.18), Inches(4.84), Inches(3.7), Inches(0.52), sz=14, bold=True, c=GD)
        T(s, desc,  x + Inches(0.18), Inches(5.42), Inches(3.7), Inches(1.0),  sz=11, c=GR)
    ftr(s)


def s07_parallel(p):
    s = blank(p)
    hdr(s, "/avengers — Five AI Engineers Working in Parallel on One Task")

    T(s, "Tasks that take one engineer one full day, completed in under 2 hours:",
      Inches(0.4), Inches(0.82), Inches(8.5), Inches(0.38), sz=13, c=GR)

    # Fury orchestrator
    R(s, Inches(0.4), Inches(1.28), Inches(8.2), Inches(0.74), fill=CHR)
    T(s, "Nick Fury  (Opus — Orchestrator)",
      Inches(0.6), Inches(1.36), Inches(4.6), Inches(0.4), sz=14, bold=True, c=W)
    T(s, "routes work  ·  validates gates  ·  blocks bad code  ·  shuts down on completion",
      Inches(5.2), Inches(1.4), Inches(3.2), Inches(0.36), sz=10, it=True, c=G)

    agents = [
        ("Planner",  "Researches codebase\n→ implementation plan", GD),
        ("Coder ×N", "Parallel code writing\nverifies build passes",G),
        ("Reviewer", "Code standards +\nOWASP security check",    G),
        ("Validator","Tests + impact analysis\nbefore any merge",   GD),
    ]
    for i, (name, desc, color) in enumerate(agents):
        x = Inches(0.4 + i * 2.06)
        R(s, x, Inches(2.2), Inches(1.88), Inches(1.62), fill=color)
        T(s, name, x + Inches(0.1), Inches(2.3),  Inches(1.68), Inches(0.4), sz=12, bold=True, c=W, a=PP_ALIGN.CENTER)
        T(s, desc, x + Inches(0.1), Inches(2.75), Inches(1.68), Inches(0.88), sz=10, c=RGBColor(0xCC, 0xFF, 0xCC), a=PP_ALIGN.CENTER)

    # Right: hero metric
    card(s, Inches(8.82), Inches(1.06), Inches(4.12), Inches(5.62), fill=LG)
    T(s, "1 DAY",   Inches(8.82), Inches(1.5),  Inches(4.12), Inches(0.9), sz=54, bold=True, c=RED, a=PP_ALIGN.CENTER)
    T(s, "↓",       Inches(8.82), Inches(2.48), Inches(4.12), Inches(0.5), sz=28, bold=True, c=G,   a=PP_ALIGN.CENTER)
    T(s, "2 HOURS", Inches(8.82), Inches(3.04), Inches(4.12), Inches(0.9), sz=54, bold=True, c=G,   a=PP_ALIGN.CENTER)
    T(s, "Airflow DAG + dbt model\n+ NR alert — all in parallel",
      Inches(8.92), Inches(4.02), Inches(3.92), Inches(0.84), sz=12, c=GR, a=PP_ALIGN.CENTER)
    T(s, "Engineer reviews\n1 PR instead of 5 files",
      Inches(8.92), Inches(4.94), Inches(3.92), Inches(0.62), sz=13, bold=True, c=GD, a=PP_ALIGN.CENTER)

    # Steps
    steps = [
        "Fury reads spec → spawns Planner + Coders simultaneously",
        "Each agent works on an independent file — zero waiting",
        "Reviewer + Validator run in parallel on all output",
        "Human sees a single clean PR to approve",
    ]
    for i, step in enumerate(steps):
        fill_ = RGBColor(0xF0, 0xFB, 0xF5) if i % 2 == 0 else W
        R(s, Inches(0.4), Inches(4.06 + i * 0.52), Inches(8.2), Inches(0.48), fill=fill_, ln=MG, lw=Pt(0.5))
        T(s, f"{i+1}.", Inches(0.56), Inches(4.14 + i * 0.52), Inches(0.34), Inches(0.36), sz=12, bold=True, c=G)
        T(s, step, Inches(0.96), Inches(4.14 + i * 0.52), Inches(7.4), Inches(0.36), sz=11, c=BK)
    ftr(s)


def s08_investigate(p):
    s = blank(p)
    hdr(s, "Incident Investigation  —  AI Queries New Relic Directly, Instantly")

    # Before
    card(s, Inches(0.4), Inches(0.82), Inches(5.8), Inches(5.55))
    top_bar(s, Inches(0.4), Inches(0.82), Inches(5.8), color=RED)
    T(s, "BEFORE  —  Manual hunt",
      Inches(0.6), Inches(0.98), Inches(5.4), Inches(0.38), sz=14, bold=True)
    T(s, "2–3\nHOURS", Inches(0.6), Inches(1.4), Inches(5.4), Inches(1.22),
      sz=56, bold=True, c=RED, a=PP_ALIGN.CENTER)
    T(s, "Log in to NR UI  ·  navigate to entity  ·  scroll logs\nswitch tabs  ·  compare charts  ·  Slack teammates",
      Inches(0.6), Inches(2.76), Inches(5.4), Inches(0.72), sz=11, c=GR, a=PP_ALIGN.CENTER)
    HL(s, Inches(0.6), Inches(3.58), Inches(5.4))
    T(s, "Form hypothesis → test → repeat →  still uncertain\nCustomer impact growing every minute",
      Inches(0.6), Inches(3.68), Inches(5.4), Inches(0.62), sz=12, c=RED, a=PP_ALIGN.CENTER)

    T(s, "→", Inches(6.38), Inches(2.9), Inches(0.55), Inches(0.6), sz=32, bold=True, c=G, a=PP_ALIGN.CENTER)

    # After
    card(s, Inches(7.1), Inches(0.82), Inches(5.8), Inches(5.55), border=G,
         fill=RGBColor(0xF0, 0xFB, 0xF5))
    top_bar(s, Inches(7.1), Inches(0.82), Inches(5.8), color=G)
    T(s, "AFTER  —  AI-powered diagnosis",
      Inches(7.3), Inches(0.98), Inches(5.4), Inches(0.38), sz=14, bold=True, c=GD)
    T(s, "30\nMINS", Inches(7.3), Inches(1.4), Inches(5.4), Inches(1.22),
      sz=56, bold=True, c=G, a=PP_ALIGN.CENTER)
    T(s, "\"analyze kafka consumer lag last 4 hours\"\nClaude queries NR, returns structured analysis",
      Inches(7.3), Inches(2.76), Inches(5.4), Inches(0.72), sz=11, c=GD, a=PP_ALIGN.CENTER)
    HL(s, Inches(7.3), Inches(3.58), Inches(5.4), G)
    T(s, "Consumer lag  ·  partition skew  ·  throughput drop\nAffected services identified — fix scoped immediately",
      Inches(7.3), Inches(3.68), Inches(5.4), Inches(0.62), sz=12, c=GD, a=PP_ALIGN.CENTER)
    T(s, "12 New Relic tools  ·  Jira auto-updated  ·  root cause in minutes",
      Inches(7.3), Inches(4.48), Inches(5.4), Inches(0.38), sz=11, bold=True, c=GD, a=PP_ALIGN.CENTER)
    ftr(s)


def s09_ship(p):
    s = blank(p)
    hdr(s, "Impact Analysis  —  Know What Breaks Before Writing a Single Line")

    R(s, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.84), fill=G)
    T(s, "We see the full impact of every change  BEFORE  it touches production.",
      Inches(0.6), Inches(0.94), Inches(12.1), Inches(0.62),
      sz=23, bold=True, c=W, a=PP_ALIGN.CENTER)

    card(s, Inches(0.4), Inches(1.84), Inches(5.9), Inches(2.44))
    top_bar(s, Inches(0.4), Inches(1.84), Inches(5.9), color=RED)
    T(s, "⚠  Before: Invisible blast radius",
      Inches(0.6), Inches(1.96), Inches(5.5), Inches(0.4), sz=14, bold=True, c=RED)
    T(s, "Engineer changes spark_ingest.py.\nDeploys. Downstream Kafka consumer breaks silently.\n5-day incident. ~$50k in eng-hours to resolve.",
      Inches(0.6), Inches(2.44), Inches(5.5), Inches(1.6), sz=12, c=GR)

    card(s, Inches(7.0), Inches(1.84), Inches(5.9), Inches(2.44), border=G,
         fill=RGBColor(0xF0, 0xFB, 0xF5))
    top_bar(s, Inches(7.0), Inches(1.84), Inches(5.9), color=G)
    T(s, "✅  After: Instant clarity",
      Inches(7.2), Inches(1.96), Inches(5.5), Inches(0.4), sz=14, bold=True, c=GD)
    T(s, "Impact analysis runs before coding.\nSees 3 callers, 1 downstream consumer.\nFix scoped. Ships confidently. Zero incident.",
      Inches(7.2), Inches(2.44), Inches(5.5), Inches(1.6), sz=12, c=BK)

    T(s, "→", Inches(6.27), Inches(2.66), Inches(0.55), Inches(0.5), sz=28, bold=True, c=G, a=PP_ALIGN.CENTER)

    metrics = [
        ("~1 sprint/qtr",   "absorbed in surprise\nrework — before",     RED),
        ("< 1 day",         "rework after impact-\naware development",    G),
        ("0 blind deploys",  "every change impact-\nchecked first",       GD),
        ("Auto-updates",    "dependency map rebuilt\non every file save",  GD),
    ]
    for i, (val, label, c_) in enumerate(metrics):
        x = Inches(0.4 + i * 3.22)
        card(s, x, Inches(4.52), Inches(3.06), Inches(1.76))
        top_bar(s, x, Inches(4.52), Inches(3.06), color=c_)
        T(s, val,   x + Inches(0.15), Inches(4.64), Inches(2.76), Inches(0.6), sz=22, bold=True, c=c_)
        T(s, label, x + Inches(0.15), Inches(5.28), Inches(2.76), Inches(0.86), sz=11, c=GR)
    ftr(s)


def s10_document(p):
    s = blank(p)
    hdr(s, "Documentation That Writes Itself  —  Jira & Confluence Always Up to Date")

    flow = [
        ("Work session\ncompletes",    W,    BK),
        ("/wrap-up\nfires (~30s)",     G,    W),
        ("Claude drafts\nsummary",     G,    W),
        ("Jira ticket\nauto-updated",  BLU,  W),
        ("Confluence\npublished",      TEAL, W),
    ]
    for i, (label, fill, tc) in enumerate(flow):
        x = Inches(0.4 + i * 2.58)
        R(s, x, Inches(0.82), Inches(2.35), Inches(1.22),
          fill=fill, ln=(G if fill == G else BLU if fill == BLU else TEAL if fill == TEAL else MG), lw=Pt(1.2))
        T(s, label, x + Inches(0.1), Inches(1.02), Inches(2.15), Inches(0.82),
          sz=12, bold=(fill != W), c=tc, a=PP_ALIGN.CENTER)
        if i < 4:
            T(s, "→", Inches(2.8 + i * 2.58), Inches(1.3), Inches(0.36), Inches(0.38),
              sz=22, bold=True, c=G, a=PP_ALIGN.CENTER)

    jira = [("Status",    "In Review → Done  ·  automatically"),
            ("Comment",   "What was done + decisions made"),
            ("Worklog",   "Time logged from session duration"),
            ("Links",     "Commits linked via branch name")]
    conf  = [("Architecture","Decisions logged after every session"),
             ("Postmortem", "Incident page auto-created on wrap-up"),
             ("Release",    "Notes from git history — auto-versioned"),
             ("Runbooks",   "Updated as team learns new patterns")]

    for col_data, label, x, acc in [
        (jira, "Jira  —  What gets auto-updated",    Inches(0.4),  BLU),
        (conf, "Confluence  —  What gets published", Inches(6.85), TEAL),
    ]:
        card(s, x, Inches(2.22), Inches(6.0), Inches(4.02))
        R(s, x, Inches(2.22), Inches(6.0), Inches(0.44), fill=acc)
        T(s, label, x + Inches(0.15), Inches(2.3), Inches(5.7), Inches(0.32), sz=14, bold=True, c=W)
        for i, (field, detail) in enumerate(col_data):
            y = Inches(2.82 + i * 0.68)
            R(s, x + Inches(0.15), y, Inches(5.7), Inches(0.62),
              fill=LG if i % 2 == 0 else W)
            T(s, field,  x + Inches(0.28), y + Inches(0.12), Inches(1.5),  Inches(0.38), sz=12, bold=True, c=acc)
            T(s, detail, x + Inches(1.9),  y + Inches(0.12), Inches(3.85), Inches(0.38), sz=11, c=BK)

    HL(s, Inches(0.4), Inches(6.4), Inches(12.5))
    T(s, "Zero manual effort.  Every session documented.  Every time.",
      Inches(0.4), Inches(6.48), Inches(12.5), Inches(0.38),
      sz=15, bold=True, c=BK, a=PP_ALIGN.CENTER)
    ftr(s)


def s11_compounds(p):
    s = blank(p)
    R(s, 0, 0, SW, SH, fill=NVY)
    R(s, 0, 0, SW, Inches(0.06), fill=G)

    T(s, "Knowledge That Compounds",
      Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
      sz=36, bold=True, c=W, a=PP_ALIGN.CENTER)
    T(s, "Every session makes every future session smarter — knowledge never leaves the team",
      Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.44),
      sz=16, it=True, c=G, a=PP_ALIGN.CENTER)

    layers = [
        ("ALWAYS LOADED",   "hot.md  — ~2k tokens, auto-loads every session start",
         "Active priorities + recent lessons + architecture summary. /wrap-up regenerates it. Zero manual effort.",
         G),
        ("PULL ON DEMAND",  "lessons.md  ·  architecture.md  ·  history.md  ·  plans/",
         "Team knowledge retrieved only when relevant — not loaded upfront. Cheaper, faster, no stale context drowning.",
         GD),
        ("REPLAYABLE WINS", "/golden save  →  /replay  →  50× faster on repeat problems",
         "After each validated fix: distill it into a golden. Next time the problem recurs: /replay loads the proven steps.",
         PUR),
    ]

    for i, (tier, title, desc, color) in enumerate(layers):
        y = Inches(2.1 + i * 1.46)
        R(s, Inches(0.4), y, Inches(12.5), Inches(1.36), fill=RGBColor(0x14, 0x14, 0x30), ln=color, lw=Pt(1.2))
        R(s, Inches(0.4), y, Inches(0.07), Inches(1.36), fill=color)
        R(s, Inches(0.56), y + Inches(0.1), Inches(1.65), Inches(0.28), fill=color)
        T(s, tier,  Inches(0.58), y + Inches(0.1),  Inches(1.61), Inches(0.26), sz=8,  bold=True, c=W, a=PP_ALIGN.CENTER)
        T(s, title, Inches(2.4),  y + Inches(0.1),  Inches(10.28), Inches(0.4), sz=15, bold=True, c=W)
        T(s, desc,  Inches(2.4),  y + Inches(0.58), Inches(10.28), Inches(0.7), sz=11, c=GR)

    R(s, Inches(0.4), Inches(6.48), Inches(12.5), Inches(0.56), fill=RGBColor(0x0A, 0x2A, 0x18))
    T(s, "Month 1: 8 issues per task  →  Month 3: 1–2 issues.  The team gets sharper, not staler.",
      Inches(0.6), Inches(6.58), Inches(12.1), Inches(0.38),
      sz=14, bold=True, c=G, a=PP_ALIGN.CENTER)


def s12_roi(p):
    """Cost & ROI — what this costs vs what it saves."""
    s = blank(p)
    hdr(s, "Return on Investment  —  Cost vs. Value, Per Engineer Per Month")

    # Cost column
    card(s, Inches(0.4), Inches(0.82), Inches(4.1), Inches(5.62))
    R(s, Inches(0.4), Inches(0.82), Inches(4.1), Inches(0.5), fill=CHR)
    T(s, "What It Costs", Inches(0.6), Inches(0.9), Inches(3.7), Inches(0.36), sz=16, bold=True, c=W)
    cost_rows = [
        ("Claude Code Pro",     "$20 / month"),
        ("Session AI spend",    "~$5–15 / day"),
        ("Monthly (est.)",      "$150–300 / eng"),
        ("Setup time",          "15 minutes"),
        ("Per-repo bootstrap",  "15 seconds"),
        ("Daily habit",         "/wrap-up 30 sec"),
    ]
    for i, (label, val) in enumerate(cost_rows):
        y = Inches(1.5 + i * 0.66)
        R(s, Inches(0.55), y, Inches(3.8), Inches(0.58), fill=LG if i % 2 == 0 else W)
        T(s, label, Inches(0.7),  y + Inches(0.12), Inches(2.0), Inches(0.36), sz=11, bold=True)
        T(s, val,   Inches(2.75), y + Inches(0.12), Inches(1.5), Inches(0.36), sz=13, bold=True, c=GD, a=PP_ALIGN.RIGHT)

    # Value column
    card(s, Inches(4.72), Inches(0.82), Inches(4.1), Inches(5.62), border=G,
         fill=RGBColor(0xF0, 0xFB, 0xF5))
    R(s, Inches(4.72), Inches(0.82), Inches(4.1), Inches(0.5), fill=G)
    T(s, "What It Saves", Inches(4.92), Inches(0.9), Inches(3.7), Inches(0.36), sz=16, bold=True, c=W)
    val_rows = [
        ("Hours saved / week",    "~10h / eng"),
        ("Eng hourly rate",       "$150 fully-loaded"),
        ("Weekly savings",        "$1,500 / eng"),
        ("Monthly savings",       "$6,000 / eng"),
        ("MTTR reduction",        "2.5h × 4/qtr saved"),
        ("Onboarding saving",     "~$5k per new hire"),
    ]
    for i, (label, val) in enumerate(val_rows):
        y = Inches(1.5 + i * 0.66)
        R(s, Inches(4.87), y, Inches(3.8), Inches(0.58), fill=RGBColor(0xE8, 0xF8, 0xF0) if i % 2 == 0 else W)
        T(s, label, Inches(5.02), y + Inches(0.12), Inches(2.0), Inches(0.36), sz=11, bold=True)
        T(s, val,   Inches(7.07), y + Inches(0.12), Inches(1.5), Inches(0.36), sz=13, bold=True, c=GD, a=PP_ALIGN.RIGHT)

    # ROI hero
    card(s, Inches(9.05), Inches(0.82), Inches(3.85), Inches(5.62), fill=G)
    T(s, "ROI",    Inches(9.05), Inches(1.55), Inches(3.85), Inches(0.55), sz=22, bold=True, c=W, a=PP_ALIGN.CENTER)
    T(s, "20–40×", Inches(9.05), Inches(2.15), Inches(3.85), Inches(1.3),  sz=64, bold=True, c=W, a=PP_ALIGN.CENTER)
    T(s, "return on\ninvestment", Inches(9.05), Inches(3.52), Inches(3.85), Inches(0.72), sz=18, c=RGBColor(0xCC, 0xFF, 0xCC), a=PP_ALIGN.CENTER)
    HL(s, Inches(9.25), Inches(4.34), Inches(3.45), W)
    T(s, "Payback period:\nfirst day of use",
      Inches(9.05), Inches(4.46), Inches(3.85), Inches(0.72), sz=14, bold=True, c=W, a=PP_ALIGN.CENTER)
    T(s, "$150–300/month\nvs $6,000 saved",
      Inches(9.05), Inches(5.22), Inches(3.85), Inches(0.72), sz=13, c=RGBColor(0xCC, 0xFF, 0xCC), a=PP_ALIGN.CENTER)

    ftr(s)


def s13_numbers(p):
    s = blank(p)
    hdr(s, "The Results  —  Numbers the Team Has Measured")

    # All numbers confirmed directly by engineer
    big4 = [
        ("8–12×",  "overall productivity\nfor DE tasks",          G),
        ("2–6h",   "dev ticket\nvs 2–3 days before",              GB),
        ("0–20m",  "Kafka/NR diagnosis\nvs 1–10 hours before",    BLU),
        ("25+ SP", "sprint capacity\n+ adhoc on top",             GD),
    ]
    for i, (num, label, c_) in enumerate(big4):
        x = Inches(0.4 + i * 3.22)
        card(s, x, Inches(0.82), Inches(3.06), Inches(3.52))
        top_bar(s, x, Inches(0.82), Inches(3.06), color=c_)
        T(s, num,   x, Inches(1.04), Inches(3.06), Inches(1.42), sz=76, bold=True, c=c_, a=PP_ALIGN.CENTER)
        T(s, label, x + Inches(0.15), Inches(2.52), Inches(2.76), Inches(0.78), sz=13, c=GR, a=PP_ALIGN.CENTER)

    # Secondary stats
    R(s, Inches(0.4), Inches(4.52), Inches(12.5), Inches(1.32), fill=LG, ln=MG, lw=Pt(0.75))
    T(s, "More measurements", Inches(0.6), Inches(4.62), Inches(12.1), Inches(0.32), sz=11, bold=True, c=GR)
    # Secondary stats — all engineer-confirmed
    secondary = [
        ("5–15 min",           "context reset after absence\n(vs hours/days before)"),
        ("3–5h",               "new pipeline (full cycle)\nvs 2–4 days before"),
        ("4–6h → 20m–1h",     "spec & plan phase"),
        ("2–4d → 1–2d",       "architecture + CDD\n+ Confluence"),
    ]
    for i, (val, label) in enumerate(secondary):
        x = Inches(0.6 + i * 3.15)
        T(s, val,   x, Inches(5.02), Inches(3.0), Inches(0.44), sz=16, bold=True, c=G)
        T(s, label, x, Inches(5.48), Inches(3.0), Inches(0.28), sz=10, c=GR)

    R(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.6), fill=G)
    T(s, "8–12× overall for DE tasks  ·  All metrics reported directly by Ingestion Team engineers",
      Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.42),
      sz=15, bold=True, c=W, a=PP_ALIGN.CENTER)
    ftr(s)


def s14_comparison(p):
    s = blank(p)
    hdr(s, "Before & After  —  Ingestion Team  ·  3 Repos  ·  3 Months")

    cols = ["What We Measured",       "Before",                                "After (AI-Led)"]
    cxs  = [Inches(0.4),              Inches(4.2),                             Inches(8.15)]
    cws  = [Inches(3.65),             Inches(3.7),                             Inches(4.9)]

    for j, (h, x, w) in enumerate(zip(cols, cxs, cws)):
        R(s, x, Inches(0.82), w, Inches(0.42),
          fill=(G if j == 2 else CHR if j == 1 else LG))
        T(s, h, x + Inches(0.12), Inches(0.9), w - Inches(0.24), Inches(0.3),
          sz=12, bold=True, c=W if j < 2 else W, a=PP_ALIGN.CENTER)

    # ✅ = confirmed directly by engineer   est. = reasonable estimate, not measured
    data = [
        ("Spec & Plan",               "4–6 hours",                "20 min – 1 hour  ✅"),
        ("Architecture + CDD",        "2–4 days",                 "1–2 days  ✅"),
        ("Development ticket",        "2–3 days",                 "2–6 hours (task-dependent)  ✅"),
        ("New pipeline (full cycle)", "2–4 days",                 "3–5 hours  ✅"),
        ("Kafka / NR diagnosis",      "1–10 hours",               "0–20 minutes (NR MCP)  ✅"),
        ("Context after absence",     "Hours / days",             "5–15 min via /bootstrap  ✅"),
        ("Sprint throughput",         "25 SP (at capacity)",      "25 SP + additional + adhoc  ✅"),
        ("Documentation effort",      "Manual, often skipped",    "Automatic via /wrap-up  ✅"),
        ("Code review layers",        "Manual review only",       "CC Opus + Copilot + Natasha + Human  ✅"),
    ]

    y = Inches(1.32)
    for i, (metric, before, after) in enumerate(data):
        fill = LG if i % 2 == 0 else W
        for j, (val, x, w) in enumerate(zip([metric, before, after], cxs, cws)):
            R(s, x, y, w, Inches(0.5), fill=fill)
            c_ = BK if j == 0 else (GR if j == 1 else GD)
            T(s, val, x + Inches(0.12), y + Inches(0.07), w - Inches(0.24), Inches(0.38),
              sz=10.5, bold=(j == 0 or j == 2), c=c_)
        y += Inches(0.5)

    R(s, Inches(0.4), y + Inches(0.06), Inches(12.6), Inches(0.36), fill=W, ln=MG, lw=Pt(0.75))
    T(s, "✅ = confirmed by Ingestion Team engineers  ·  Repos: om-airflow-dags  ·  spark-kafka-apps  ·  tf-dataos-new-relic",
      Inches(0.6), y + Inches(0.12), Inches(12.2), Inches(0.28), sz=10, c=GR, a=PP_ALIGN.CENTER)
    ftr(s)


def s15_control(p):
    s = blank(p)
    hdr(s, "AI Does the Work.  Humans Stay in Control.  Always.")

    gates = [
        ("1", "Requirements",     "Human writes the spec — AI reads it, asks clarifying questions"),
        ("2", "Architecture",     "Human approves the implementation plan — before a line is written"),
        ("3", "Security Review",  "AI flags every OWASP finding — any critical issue blocks the pipeline"),
        ("4", "PR Merge",         "Human reviews every PR — agent output treated like any other code"),
        ("5", "Deploy",           "Human initiates the deploy — AI writes release notes + monitors"),
    ]
    y = Inches(0.82)
    for num, title, detail in gates:
        is3 = num == "3"
        h = Inches(1.0)
        R(s, Inches(0.4), y, Inches(12.5), h,
          fill=(RGBColor(0xF0, 0xFB, 0xF5) if is3 else W),
          ln=(G if is3 else MG), lw=Pt(1.5 if is3 else 0.75))
        R(s, Inches(0.4), y, Inches(0.54), h, fill=G if is3 else CHR)
        T(s, num,    Inches(0.4),   y + Inches(0.27), Inches(0.54), Inches(0.46),
          sz=18, bold=True, c=W, a=PP_ALIGN.CENTER)
        T(s, title,  Inches(1.1),   y + Inches(0.12), Inches(3.8), Inches(0.42), sz=15, bold=True)
        T(s, detail, Inches(5.06),  y + Inches(0.2),  Inches(7.7), Inches(0.6),  sz=13, c=GR)
        y += h + Inches(0.1)

    R(s, Inches(0.4), y + Inches(0.06), Inches(12.5), Inches(0.38), fill=LG, ln=MG)
    T(s, "Agent PRs treated identically to human PRs — same review, same standards, same gates.",
      Inches(0.6), y + Inches(0.12), Inches(12.1), Inches(0.28), sz=11, it=True, c=GR, a=PP_ALIGN.CENTER)
    ftr(s)


def s16_start(p):
    s = blank(p)
    hdr(s, "Getting Started  —  From Zero to Productive in 4 Steps")

    steps = [
        ("15 MIN",   "Install",      "One-time per engineer",       "install.sh checks prereqs, installs 20 skills,\nhooks, MCPs, budget dial — fully automated.", G),
        ("15 SEC",   "/bootstrap",   "First visit to any repo",     "Detects stack, seeds memory, builds code graph,\nwires all MCPs — runs once, lasts forever.", GD),
        ("ALWAYS",   "Work",         "Skills auto-trigger",         "20 skills activate on file type or topic.\nCopilot handles lines. Claude handles investigations.", G),
        ("30 SEC",   "/wrap-up",     "Every session end",           "Persists history, updates Jira, refreshes graph,\nrebuilds hot.md. Habit that compounds daily.", GD),
    ]
    for i, (time, cmd, when, desc, color) in enumerate(steps):
        x = Inches(0.4 + i * 3.22)
        card(s, x, Inches(0.82), Inches(3.06), Inches(4.7))
        R(s, x, Inches(0.82), Inches(3.06), Inches(0.72), fill=color)
        T(s, time, x + Inches(0.1), Inches(0.9),  Inches(2.86), Inches(0.54), sz=18, bold=True, c=W, a=PP_ALIGN.CENTER)
        T(s, f"Step {i+1}  —  {cmd}", x + Inches(0.14), Inches(1.68), Inches(2.78), Inches(0.44), sz=14, bold=True, c=color)
        T(s, when,  x + Inches(0.14), Inches(2.18), Inches(2.78), Inches(0.32), sz=11, it=True, c=GR)
        T(s, desc,  x + Inches(0.14), Inches(2.58), Inches(2.78), Inches(1.7),  sz=11, c=BK)

    HL(s, Inches(0.4), Inches(5.7), Inches(12.5))
    T(s, "Week-by-week rollout:", Inches(0.4), Inches(5.8), Inches(3.0), Inches(0.3), sz=11, bold=True, c=GR)
    weeks = [
        ("Week 1", "Install + /bootstrap 1 repo + /wrap-up daily"),
        ("Week 2", "All repos + first /golden save after a win"),
        ("Week 3", "/avengers on a real multi-file task"),
        ("Week 4+","Team golden library grows · 15–25× steady"),
    ]
    for i, (w, desc) in enumerate(weeks):
        x = Inches(0.4 + i * 3.22)
        T(s, w,    x, Inches(6.16), Inches(3.0), Inches(0.3), sz=11, bold=True, c=G)
        T(s, desc, x, Inches(6.5),  Inches(3.0), Inches(0.62), sz=10, c=GR)
    ftr(s)


def s17_ask(p):
    s = blank(p)
    R(s, 0, 0, SW, SH, fill=NVY)
    R(s, 0, 0, SW, Inches(0.06), fill=G)
    R(s, Inches(10.2), 0, Inches(3.13), SH, fill=RGBColor(0x12, 0x12, 0x30))

    T(s, "⬢ new relic", Inches(0.6), Inches(0.3), Inches(3), Inches(0.5), sz=20, bold=True, c=G)
    T(s, "What We're Asking",
      Inches(0.6), Inches(0.98), Inches(9.3), Inches(0.66), sz=34, bold=True, c=W)
    HL(s, Inches(0.6), Inches(1.7), Inches(8.0), G)

    asks = [
        ("Support adoption",  "Drive AI-led workflows across all Data Engineering & Management in DTM."),
        ("Allocate time",     "Engineers need 30 min/day: /wrap-up + goldens. The compound returns are real."),
        ("Endorse the model", "Establish this as the standard template for all DE teams at New Relic."),
    ]
    for i, (label, detail) in enumerate(asks):
        y = Inches(2.0 + i * 1.18)
        R(s, Inches(0.6), y, Inches(9.38), Inches(1.02), fill=RGBColor(0x14, 0x14, 0x38), ln=G, lw=Pt(0.75))
        R(s, Inches(0.6), y, Inches(0.07), Inches(1.02), fill=G)
        T(s, label,  Inches(0.84), y + Inches(0.08), Inches(2.5),  Inches(0.4), sz=16, bold=True, c=G)
        T(s, detail, Inches(3.45), y + Inches(0.14), Inches(6.38), Inches(0.72), sz=14, c=W)

    # Right panel numbers
    for i, (val, lbl) in enumerate([("15–25×","sustained avg"),("10×","novel work"),("50×","repeat tasks"),("30 min","to root cause"),("20–40×","ROI")]):
        T(s, val, Inches(10.2), Inches(1.3 + i * 1.08), Inches(3.13), Inches(0.68), sz=26, bold=True, c=G, a=PP_ALIGN.CENTER)
        T(s, lbl, Inches(10.2), Inches(1.95 + i * 1.08), Inches(3.13), Inches(0.3), sz=10, it=True, c=GR, a=PP_ALIGN.CENTER)

    T(s, "Thank You  ·  Questions?",
      Inches(0.6), Inches(5.78), Inches(9.3), Inches(0.65), sz=28, c=GR)
    T(s, "Repo: github.com/nr-mdakilahmed/claude_code_setup",
      Inches(0.6), Inches(6.5), Inches(9.3), Inches(0.38), sz=12, it=True, c=RGBColor(0x44, 0x44, 0x66))


# ══════════════════════════════════════════════════════════════════════════════

SLIDES = [
    # Act 1 — Story: Challenge
    s01_title,
    s02_problem,
    # Act 2 — Our approach
    s03_how_we_work,       # DataOS Ingestion: how we work now
    s03b_two_tools,        # Copilot + Claude Code: two tools, different jobs
    s04_architecture,      # System architecture diagram
    # Act 3 — The lifecycle
    s05_lifecycle,         # Full lifecycle with verified timings
    # Act 4 — Benefits (capabilities in action)
    s06_onboard,           # Onboard in hours
    s07_parallel,          # /avengers parallel build
    s08_investigate,       # NR MCP diagnosis
    s09_ship,              # Impact analysis
    s10_document,          # Auto-documentation
    s11_compounds,         # Knowledge compounds
    # Act 5 — Proof
    s13_numbers,           # Real metrics
    s14_comparison,        # Before & after
    # Act 6 — Governance + Getting started + Ask
    s15_control,           # AI works, humans in control
    s16_start,             # Getting started
    s17_ask,               # What we're asking
]

OUT_PPTX = os.path.join(DIR, "Ingestion-Team-AI-Led-Demo.pptx")
OUT_PDF  = os.path.join(DIR, "Ingestion-Team-AI-Led-Demo.pdf")

def main():
    pr = prs()
    for fn in SLIDES:
        fn(pr)
    pr.save(OUT_PPTX)
    print(f"✓  {len(pr.slides)} slides  →  {OUT_PPTX}")

    soffice = "/opt/homebrew/bin/soffice"
    if os.path.exists(soffice):
        r = subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                            "--outdir", DIR, OUT_PPTX], capture_output=True, text=True)
        print(f"✓  PDF  →  {OUT_PDF}" if r.returncode == 0 else f"⚠  PDF: {r.stderr.strip()}")

if __name__ == "__main__":
    main()
